from pathlib import Path
import json
from pptx import Presentation

DATA_DIR = Path("data")
OUT_PATH = Path("slides.jsonl")

def iter_pptx_slides(pptx_path: Path):
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides, start=1):
        text_runs = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text.strip())
        text = "\n".join(t for t in text_runs if t)
        yield {
            "doc_name": pptx_path.stem,
            "path": str(pptx_path.resolve()),
            "slide": i,
            "text": text
        }

if __name__ == "__main__":
    with OUT_PATH.open("w", encoding="utf-8") as f:
        for pptx_path in sorted(DATA_DIR.glob("*.pptx")):
            for rec in iter_pptx_slides(pptx_path):
                rec["char_count"] = len(rec["text"])
                f.write(json.dumps(rec, ensure_ascii=False) + "\n")
    print(f"✅ 完了！ {OUT_PATH} に保存しました")

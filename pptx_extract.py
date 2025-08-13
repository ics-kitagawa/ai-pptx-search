# pptx_extract.py
from __future__ import annotations
from pathlib import Path
import json
from typing import List
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

DATA_DIR = Path("data")
OUT_PATH = Path("slides.jsonl")

def extract_textframe(shape) -> List[str]:
    """箇条書きの階層を '- ' とインデントで素直に表現"""
    lines = []
    if not shape.has_text_frame:
        return lines
    tf = shape.text_frame
    for p in tf.paragraphs:
        # runsが空でもp.textに段落全体のテキストが入っている
        text = "".join(r.text for r in p.runs).strip() or (p.text or "").strip()
        if not text:
            continue
        indent = "  " * getattr(p, "level", 0)
        lines.append(f"{indent}- {text}")
    return lines

def extract_table(shape) -> str:
    """表を Markdown ふうに直列化（検索しやすくする）"""
    tbl = shape.table
    rows = []
    for r in tbl.rows:
        cells = []
        for c in r.cells:
            # セル内改行は空白に寄せる（お好みで '\n' のままでも可）
            cells.append(" ".join((c.text or "").split()))
        rows.append(cells)

    if not rows:
        return ""

    # ヘッダー推定：一律で1行目をヘッダー扱い（必要なら自動判定に変えてOK）
    header = rows[0]
    body = rows[1:] if len(rows) > 1 else []

    # Markdown 風
    out = []
    out.append("| " + " | ".join(header) + " |")
    out.append("| " + " | ".join(["---"] * len(header)) + " |")
    for r in body:
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out)

def walk_shape(shape) -> List[str]:
    """図形を再帰的に歩いてテキストを収集"""
    blocks: List[str] = []
    st = shape.shape_type

    # グループ化図形
    if st == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
        for ch in shape.shapes:
            blocks.extend(walk_shape(ch))
        return blocks

    # 表
    if getattr(shape, "has_table", False):
        t = extract_table(shape)
        if t.strip():
            blocks.append(t)
        return blocks

    # テキストフレーム
    if getattr(shape, "has_text_frame", False):
        lines = extract_textframe(shape)
        if lines:
            blocks.append("\n".join(lines))
        return blocks

    # その他（図やチャートはテキストが取れないことが多い）
    return blocks

def extract_slide_text(slide) -> str:
    # 図形本体
    blocks: List[str] = []
    for shape in slide.shapes:
        blocks.extend(walk_shape(shape))

    # ノート（必要ならON）
    if getattr(slide, "has_notes_slide", False) and slide.notes_slide:
        notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        if notes:
            blocks.append("### Notes\n" + notes)

    # ブロック間は空行で区切る（表もテキストも区別しやすく）
    return "\n\n".join([b for b in blocks if b.strip()])

def iter_pptx_slides(pptx_path: Path):
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides, start=1):
        text = extract_slide_text(slide)
        yield {
            "doc_name": pptx_path.stem,
            "path": str(pptx_path.resolve()),
            "slide": i,
            "text": text,
        }

def main():
    out = OUT_PATH.open("w", encoding="utf-8")
    try:
        # サブフォルダも含めて *.pptx を拾いたい場合は rglob に変更
        for pptx_path in sorted(DATA_DIR.glob("*.pptx")):
            for rec in iter_pptx_slides(pptx_path):
                rec["char_count"] = len(rec["text"])
                out.write(json.dumps(rec, ensure_ascii=False) + "\n")
    finally:
        out.close()
    print(f"✅ 完了: {OUT_PATH} を生成しました")

if __name__ == "__main__":
    main()
